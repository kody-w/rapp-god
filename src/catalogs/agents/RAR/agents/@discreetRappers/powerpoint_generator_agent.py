"""
PowerPoint Generator Agent V2 - Template-Based Microsoft Design
Purpose: Generate professional PowerPoint presentations using Microsoft templates

Design principles:
- Template-based generation for consistent branding
- Supports multiple templates (BaseTemplateBlue, ZavaTemplate, etc.)
- Smart layout selection based on content type
- Proper placeholder population
- Fallback to programmatic generation if template not available

Templates supported:
- BaseTemplateBlue.pptx: Microsoft corporate template (113 layouts)
- ZavaTemplate.pptx: Modern business template (62 layouts)
- BaseTemplateDynamics.pptx: Dynamics-style template

Usage:
1. With template: action="create_presentation", template="BaseTemplateBlue", slides=[...]
2. Without template: action="create_presentation", slides=[...] (uses default styling)
"""

from __future__ import annotations

# ═══════════════════════════════════════════════════════════════
# RAPP AGENT MANIFEST — Do not remove. Used by registry builder.
# ═══════════════════════════════════════════════════════════════
__manifest__ = {
    "schema": "rapp-agent/1.0",
    "name": "@discreetRappers/powerpoint_generator_agent",
    "version": "1.0.2",
    "display_name": "PowerPointGeneratorV2",
    "description": "Generates PowerPoint decks from slide specs with python-pptx, using Microsoft templates and smart layout selection.",
    "author": "Bill Whalen",
    "tags": ["productivity", "powerpoint", "presentations", "templates", "microsoft"],
    "category": "productivity",
    "quality_tier": "community",
    "requires_env": [],
    "dependencies": ["@rapp/basic_agent"],
}
# ═══════════════════════════════════════════════════════════════


import json
import logging
import os
from datetime import datetime
from typing import Dict, Any, List, Optional, Tuple
from agents.basic_agent import BasicAgent
from utils.storage_factory import get_storage_manager

# Import python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.oxml.ns import nsmap
    PPTX_AVAILABLE = True
except ImportError as e:
    PPTX_AVAILABLE = False
    PPTX_IMPORT_ERROR = str(e)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PowerPointGeneratorAgentV2(BasicAgent):
    """
    Agent for generating professional presentations using Microsoft templates.
    """

    # Template configurations - maps template names to layout indexes
    TEMPLATE_CONFIGS = {
        "PowerpointTemplateBlue": {
            "file": "docs/ppt/ppt_templates/PowerpointTemplateBlue.pptx",
            "layouts": {
                "title": 2,  # "Title Slide"
                "title_photo": 0,  # "Title square photo"
                "section": 102,  # "Section Divider"
                "content": 9,  # "Title and Content"
                "two_column": 11,  # "Two Column Bullet text"
                "three_column": 15,  # "Three Column Bullet with Subtitles"
                "four_column": 16,  # "Four Column Bullet with Subtitles"
                "comparison": 11,  # "Two Column Bullet text"
                "quote": 64,  # "Quote slide 1b"
                "code": 94,  # "Developer Code Layout full page"
                "demo": 100,  # "Demo slide"
                "blank": 108,  # "Blank 12 Column"
                "closing": 111,  # "Closing logo slide"
                "title_only": 18,  # "Title Only"
            }
        },
        "BaseTemplateBlueV2": {
            "file": "docs/ppt/BaseTemplateBlueV2.pptx",
            "layouts": {
                "title": 2,  # "Title Slide"
                "title_photo": 0,  # "Title square photo"
                "section": 102,  # "Section Divider"
                "content": 9,  # "Title and Content"
                "two_column": 11,  # "Two Column Bullet text"
                "three_column": 15,  # "Three Column Bullet with Subtitles"
                "four_column": 16,  # "Four Column Bullet with Subtitles"
                "comparison": 11,  # "Two Column Bullet text"
                "quote": 64,  # "Quote slide 1b"
                "code": 94,  # "Developer Code Layout full page"
                "demo": 100,  # "Demo slide"
                "blank": 108,  # "Blank 12 Column"
                "closing": 111,  # "Closing logo slide"
                "title_only": 18,  # "Title Only"
            }
        },
        "BaseTemplateBlue": {
            "file": "docs/ppt/BaseTemplateBlue.pptx",
            "layouts": {
                "title": 2,  # "Title Slide"
                "title_photo": 0,  # "Title square photo"
                "section": 102,  # "Section Divider"
                "content": 9,  # "Title and Content"
                "two_column": 11,  # "Two Column Bullet text"
                "three_column": 15,  # "Three Column Bullet with Subtitles"
                "four_column": 16,  # "Four Column Bullet with Subtitles"
                "comparison": 11,  # "Two Column Bullet text"
                "quote": 64,  # "Quote slide 1b"
                "code": 94,  # "Developer Code Layout full page"
                "demo": 100,  # "Demo slide"
                "blank": 108,  # "Blank 12 Column"
                "closing": 111,  # "Closing logo slide"
                "title_only": 18,  # "Title Only"
            }
        },
        "ZavaTemplate": {
            "file": "docs/ppt/ZavaTemplate.pptx",
            "layouts": {
                "title": 0,  # "Title 1"
                "title_photo": 10,  # "Title Photo 1"
                "section": 14,  # "Section Header 1"
                "content": 24,  # "Content 1"
                "two_column": 41,  # "Two Content"
                "comparison": 43,  # "Comparison"
                "quote": 59,  # "Quote"
                "statement": 56,  # "Statement"
                "number": 53,  # "Number Large"
                "conclusion": 48,  # "Conclusion 1"
                "blank": 45,  # "Blank"
                "title_only": 44,  # "Title Only"
                "agenda": 20,  # "Agenda"
            }
        },
        "BaseTemplateDynamics": {
            "file": "docs/ppt/BaseTemplateDynamics.pptx",
            "layouts": {
                "title": 0,
                "content": 1,
                "blank": 6,
            }
        }
    }

    # Microsoft color palette
    COLORS = {
        "ms_blue": "0078D4",
        "ms_dark_blue": "004578",
        "ms_light_blue": "50E6FF",
        "ms_green": "107C10",
        "ms_red": "D13438",
        "ms_orange": "FF8C00",
        "ms_purple": "5C2D91",
        "black": "000000",
        "dark_gray": "323130",
        "medium_gray": "605E5C",
        "light_gray": "A19F9D",
        "white": "FFFFFF",
    }

    # Segoe UI fonts (Microsoft standard)
    FONTS = {
        "title": {"name": "Segoe UI Semibold", "size": 44, "bold": False},
        "subtitle": {"name": "Segoe UI", "size": 24, "bold": False},
        "heading": {"name": "Segoe UI Semibold", "size": 28, "bold": False},
        "body": {"name": "Segoe UI", "size": 18, "bold": False},
        "caption": {"name": "Segoe UI", "size": 14, "bold": False},
    }

    def __init__(self):
        self.name = 'PowerPointGeneratorV2'
        self.metadata = {
            "name": self.name,
            "description": """Generate professional PowerPoint presentations using Microsoft templates.

Templates available:
- BaseTemplateBlue: Microsoft corporate template (recommended)
- ZavaTemplate: Modern business template
- BaseTemplateDynamics: Dynamics-style template

Actions:
- create_presentation: Create multi-slide presentation
- list_templates: List available templates and their layouts
- list_layouts: List layouts for a specific template

Slide types: title, section, content, two_column, comparison, quote, stats, pipeline, blank

Example:
{
  "action": "create_presentation",
  "customer": "Contoso",
  "template": "BaseTemplateBlue",
  "output_filename": "my_presentation",
  "slides": [
    {"type": "title", "title": "My Presentation", "subtitle": "Subtitle here"},
    {"type": "content", "title": "Key Points", "bullets": ["Point 1", "Point 2"]},
    {"type": "comparison", "title": "Before vs After", "left_label": "Before", "right_label": "After", "left_items": ["Old way"], "right_items": ["New way"]}
  ]
}""",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "enum": ["create_presentation", "list_templates", "list_layouts"]
                    },
                    "customer": {
                        "type": "string",
                        "description": "Customer name - creates a subfolder in docs/ppt for this customer"
                    },
                    "template": {
                        "type": "string",
                        "description": "Template name (BaseTemplateBlue, ZavaTemplate, BaseTemplateDynamics)"
                    },
                    "slides": {
                        "type": "array",
                        "items": {"type": "object"}
                    },
                    "output_filename": {"type": "string"},
                },
                "required": ["action"]
            }
        }
        super().__init__(self.name, self.metadata)

        try:
            self.storage = get_storage_manager()
        except Exception as e:
            logger.warning(f"Storage not available: {e}")
            self.storage = None

        # Find base path for templates
        self.base_path = self._find_base_path()

    def _find_base_path(self) -> str:
        """Find the base path for the RAPP project."""
        # Try common locations
        possible_paths = [
            os.getcwd(),
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
            "c:/Users/billwhalen/OneDrive - Microsoft/Documents/GitHub/RAPP/CommunityRAPP-main",
        ]
        for path in possible_paths:
            if os.path.exists(os.path.join(path, "docs", "ppt")):
                return path
        return os.getcwd()

    def perform(self, **kwargs) -> str:
        """Execute the requested action."""
        if not PPTX_AVAILABLE:
            return json.dumps({
                "status": "error",
                "error": f"python-pptx library not available: {PPTX_IMPORT_ERROR}",
                "suggestion": "Install with: pip install python-pptx"
            })

        action = kwargs.get('action', 'create_presentation')

        try:
            if action == 'list_templates':
                return self._list_templates()
            elif action == 'list_layouts':
                return self._list_layouts(kwargs.get('template', 'BaseTemplateBlue'))
            elif action == 'create_presentation':
                return self._create_presentation(**kwargs)
            else:
                return json.dumps({
                    "status": "error",
                    "error": f"Unknown action: {action}",
                    "available_actions": ["create_presentation", "list_templates", "list_layouts"]
                })
        except Exception as e:
            logger.error(f"PowerPoint generation error: {e}")
            import traceback
            return json.dumps({
                "status": "error",
                "error": str(e),
                "traceback": traceback.format_exc()
            })

    def _list_templates(self) -> str:
        """List available templates."""
        templates = {}
        for name, config in self.TEMPLATE_CONFIGS.items():
            # Normalize path separators for Windows
            template_rel_path = config["file"].replace("/", os.sep)
            template_path = os.path.join(self.base_path, template_rel_path)
            templates[name] = {
                "file": config["file"],
                "exists": os.path.exists(template_path),
                "layouts": list(config["layouts"].keys())
            }
        return json.dumps({"status": "success", "templates": templates}, indent=2)

    def _list_layouts(self, template_name: str) -> str:
        """List layouts for a specific template."""
        if template_name not in self.TEMPLATE_CONFIGS:
            return json.dumps({
                "status": "error",
                "error": f"Unknown template: {template_name}",
                "available": list(self.TEMPLATE_CONFIGS.keys())
            })

        config = self.TEMPLATE_CONFIGS[template_name]
        # Normalize path separators for Windows
        template_rel_path = config["file"].replace("/", os.sep)
        template_path = os.path.join(self.base_path, template_rel_path)

        if not os.path.exists(template_path):
            return json.dumps({
                "status": "error",
                "error": f"Template file not found: {template_path}"
            })

        # Handle .potx files by converting to .pptx in temp location
        import tempfile
        import shutil
        
        actual_path = template_path
        if template_path.lower().endswith('.potx'):
            temp_dir = tempfile.gettempdir()
            temp_pptx = os.path.join(temp_dir, f"temp_template_{template_name}.pptx")
            shutil.copy2(template_path, temp_pptx)
            actual_path = temp_pptx

        prs = Presentation(actual_path)
        layouts = []
        for i, layout in enumerate(prs.slide_layouts):
            layouts.append({"index": i, "name": layout.name})

        return json.dumps({
            "status": "success",
            "template": template_name,
            "layout_count": len(layouts),
            "mapped_layouts": config["layouts"],
            "all_layouts": layouts
        }, indent=2)

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor."""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )

    def _create_presentation(self, **kwargs) -> str:
        """Create a presentation using templates."""
        template_name = kwargs.get('template', 'BaseTemplateBlue')
        slides = kwargs.get('slides', [])
        output_filename = kwargs.get('output_filename', 'presentation')

        if not slides:
            return json.dumps({
                "status": "error",
                "error": "No slides provided. Use 'slides' parameter with array of slide configs."
            })

        # Load template or create blank presentation
        prs = self._load_template(template_name)
        if prs is None:
            return json.dumps({
                "status": "error",
                "error": f"Could not load template: {template_name}"
            })

        config = self.TEMPLATE_CONFIGS.get(template_name, {})
        layout_map = config.get("layouts", {})

        # Process each slide
        for i, slide_config in enumerate(slides):
            slide_type = slide_config.get('type', 'content')
            self._add_slide(prs, slide_config, slide_type, layout_map, i + 1)

        return self._save_presentation(prs, output_filename, kwargs)

    def _remove_placeholder_shapes(self, slide) -> None:
        """Remove placeholder shapes from a slide to avoid template artifacts."""
        shapes_to_remove = []
        for shape in slide.shapes:
            # Check if it's a placeholder shape
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                shapes_to_remove.append(shape)
        
        # Remove the placeholder shapes
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    def _load_template(self, template_name: str) -> Optional[Presentation]:
        """Load a PowerPoint template."""
        if template_name not in self.TEMPLATE_CONFIGS:
            logger.warning(f"Unknown template {template_name}, using blank presentation")
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            return prs

        config = self.TEMPLATE_CONFIGS[template_name]
        # Normalize path separators for Windows
        template_rel_path = config["file"].replace("/", os.sep)
        template_path = os.path.join(self.base_path, template_rel_path)

        if not os.path.exists(template_path):
            logger.warning(f"Template file not found: {template_path}")
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            return prs

        try:
            # Handle .potx files by converting to .pptx in temp location
            import tempfile
            import shutil
            
            if template_path.lower().endswith('.potx'):
                # Copy .potx to temp .pptx file (python-pptx doesn't support .potx directly)
                temp_dir = tempfile.gettempdir()
                temp_pptx = os.path.join(temp_dir, f"temp_template_{template_name}.pptx")
                shutil.copy2(template_path, temp_pptx)
                template_path = temp_pptx
            
            prs = Presentation(template_path)
            # Remove any existing slides from template
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
            return prs
        except Exception as e:
            logger.error(f"Error loading template: {e}")
            return None

    def _add_slide(self, prs: Presentation, config: Dict, slide_type: str, 
                   layout_map: Dict, page_num: int) -> None:
        """Add a slide based on type and configuration."""
        # Get the appropriate layout
        layout_idx = layout_map.get(slide_type, layout_map.get('content', 0))

        # Ensure layout index is valid
        if layout_idx >= len(prs.slide_layouts):
            layout_idx = 0

        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        # Populate the slide based on type
        if slide_type == 'title':
            self._populate_title_slide(slide, config)
        elif slide_type == 'section':
            self._populate_section_slide(slide, config)
        elif slide_type == 'content':
            self._populate_content_slide(slide, config)
        elif slide_type in ['two_column', 'comparison']:
            self._populate_comparison_slide(slide, config)
        elif slide_type == 'quote':
            self._populate_quote_slide(slide, config)
        elif slide_type == 'stats':
            self._populate_stats_slide(slide, prs, config)
        elif slide_type == 'pipeline':
            self._populate_pipeline_slide(slide, prs, config)
        elif slide_type == 'image':
            self._populate_image_slide(slide, config)
        elif slide_type == 'title_image':
            self._populate_title_image_slide(slide, config)
        elif slide_type == 'value_cards':
            self._populate_value_cards_slide(slide, config)
        elif slide_type == 'before_after':
            self._populate_before_after_slide(slide, config)
        elif slide_type == 'agent_cards':
            self._populate_agent_cards_slide(slide, config)
        elif slide_type == 'metric_boxes':
            self._populate_metric_boxes_slide(slide, config)
        elif slide_type == 'process_flow':
            self._populate_process_flow_slide(slide, config)
        else:
            # Default content slide
            self._populate_content_slide(slide, config)

    def _populate_image_slide(self, slide, config: Dict) -> None:
        """Populate a slide with an image."""
        title = config.get('title', '')
        image_path = config.get('image_path', '')
        caption = config.get('caption', '')
        
        # Remove non-title placeholders to avoid artifacts
        shapes_to_remove = []
        title_shape = None
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                if shape.placeholder_format.type == 1:  # Title placeholder
                    title_shape = shape
                else:
                    shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
        
        # Set title
        if title_shape and title:
            title_shape.text_frame.paragraphs[0].text = title
            self._style_text(title_shape.text_frame.paragraphs[0], "heading")
        elif title:
            self._add_title_textbox(slide, title, 0.5, 0.3, 12.333, size=28)
        
        # Add image if path exists
        if image_path and os.path.exists(image_path):
            # Calculate centered position
            img_width = Inches(10)
            img_left = Inches(1.667)  # Center on 13.333" wide slide
            img_top = Inches(1.3)
            img_height = Inches(5.5)
            
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
        
        # Add caption if provided
        if caption:
            caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4))
            tf = caption_box.text_frame
            p = tf.paragraphs[0]
            p.text = caption
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["medium_gray"])
            p.alignment = PP_ALIGN.CENTER

    def _populate_title_image_slide(self, slide, config: Dict) -> None:
        """Populate a slide with title, content bullets, and an image side by side."""
        title = config.get('title', '')
        content = config.get('content', [])
        image_path = config.get('image_path', '')
        
        # Set title
        title_set = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.placeholder_format:
                if shape.placeholder_format.type == 1:
                    shape.text_frame.paragraphs[0].text = title
                    self._style_text(shape.text_frame.paragraphs[0], "heading")
                    title_set = True
                    break
        
        if not title_set and title:
            self._add_title_textbox(slide, title, 0.5, 0.3, 12.333, size=28)
        
        # Add content on left side
        if content:
            self._add_bullet_textbox(slide, content, 0.5, 1.3, 5.5, 5.5)
        
        # Add image on right side
        if image_path and os.path.exists(image_path):
            img_left = Inches(6.5)
            img_top = Inches(1.3)
            img_width = Inches(6.3)
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

    def _populate_title_slide(self, slide, config: Dict) -> None:
        """Populate a title slide."""
        # Remove all placeholder shapes to avoid template artifacts
        self._remove_placeholder_shapes(slide)
        
        title = config.get('title', '')
        subtitle = config.get('subtitle', '')

        # Title slides have dark background - use white text
        self._add_title_textbox(slide, title, 0.5, 2.5, 12.333, color="#FFFFFF")
        if subtitle:
            self._add_subtitle_textbox(slide, subtitle, 0.5, 3.5, 12.333, color="#CCCCCC")

    def _populate_section_slide(self, slide, config: Dict) -> None:
        """Populate a section divider slide."""
        # Remove all placeholder shapes to avoid template artifacts
        self._remove_placeholder_shapes(slide)
        
        title = config.get('title', '')

        # Section slides have dark background - use white text
        self._add_title_textbox(slide, title, 0.5, 3.0, 12.333, size=36, color="#FFFFFF")

    def _populate_content_slide(self, slide, config: Dict) -> None:
        """Populate a content slide with bullets."""
        # Remove all placeholder shapes to avoid template artifacts
        self._remove_placeholder_shapes(slide)
        
        title = config.get('title', '')
        bullets = config.get('bullets', config.get('content', []))

        # Add title and content as textboxes (placeholders removed)
        if title:
            self._add_title_textbox(slide, title, 0.5, 0.5, 12.333, size=28)
        if bullets:
            self._add_bullet_textbox(slide, bullets, 0.5, 1.5, 12.333, 5.5)

    def _populate_comparison_slide(self, slide, config: Dict) -> None:
        """Populate a comparison/two-column slide."""
        # Remove all placeholder shapes to avoid template artifacts
        self._remove_placeholder_shapes(slide)
        
        title = config.get('title', '')
        
        # Support both old format (left_label/right_label) and new format (left/right objects)
        left_data = config.get('left', {})
        right_data = config.get('right', {})
        
        if isinstance(left_data, dict):
            # New format with nested title/content
            left_label = left_data.get('title', config.get('left_label', 'Left'))
            left_items = left_data.get('content', config.get('left_items', []))
        else:
            left_label = config.get('left_label', 'Before')
            left_items = config.get('left_items', [])
            
        if isinstance(right_data, dict):
            right_label = right_data.get('title', config.get('right_label', 'Right'))
            right_items = right_data.get('content', config.get('right_items', []))
        else:
            right_label = config.get('right_label', 'After')
            right_items = config.get('right_items', [])

        # Add title as textbox (placeholders removed)
        if title:
            self._add_title_textbox(slide, title, 0.5, 0.3, 12.333, size=28)

        # Add comparison content via text boxes
        self._add_two_column_content(slide, left_label, right_label, left_items, right_items)

    def _populate_quote_slide(self, slide, config: Dict) -> None:
        """Populate a quote slide."""
        # Remove all placeholder shapes to avoid template artifacts
        self._remove_placeholder_shapes(slide)
        
        quote = config.get('quote', '')
        author = config.get('author', config.get('quote_author', ''))

        # Add quote box directly (placeholders removed)
        self._add_quote_box(slide, quote, author)

    def _populate_stats_slide(self, slide, prs, config: Dict) -> None:
        """Populate a stats/metrics slide."""
        # Remove all placeholder shapes to avoid template artifacts
        self._remove_placeholder_shapes(slide)
        
        title = config.get('title', '')
        stats = config.get('stats', config.get('metrics', []))

        # Add title as textbox (placeholders removed)
        if title:
            self._add_title_textbox(slide, title, 0.5, 0.3, 12.333, size=28)

        # Add stats boxes
        self._add_stats_boxes(slide, prs, stats)

    def _populate_pipeline_slide(self, slide, prs, config: Dict) -> None:
        """Populate a pipeline/process slide."""
        # Remove all placeholder shapes to avoid template artifacts
        self._remove_placeholder_shapes(slide)
        
        title = config.get('title', '')
        steps = config.get('steps', [])

        # Add title as textbox (placeholders removed)
        if title:
            self._add_title_textbox(slide, title, 0.5, 0.3, 12.333, size=28)

        # Add pipeline visualization
        self._add_pipeline_boxes(slide, prs, steps)

    def _populate_value_cards_slide(self, slide, config: Dict) -> None:
        """Populate a slide with value proposition cards (like the HTML demo)."""
        # Remove template placeholders first
        self._remove_placeholder_shapes(slide)
        
        title = config.get('title', '')
        cards = config.get('cards', [])
        
        # Set title
        self._add_title_textbox(slide, title, 0.5, 0.3, 12.333, size=28)
        
        # Calculate card positions (up to 4 cards per row)
        num_cards = len(cards)
        card_width = 3.8
        card_height = 2.8
        gap = 0.3
        
        if num_cards <= 3:
            start_x = (13.333 - (num_cards * card_width + (num_cards - 1) * gap)) / 2
            cards_per_row = num_cards
        else:
            start_x = (13.333 - (3 * card_width + 2 * gap)) / 2
            cards_per_row = 3
        
        for i, card in enumerate(cards):
            row = i // cards_per_row
            col = i % cards_per_row
            x = start_x + col * (card_width + gap)
            y = 1.3 + row * (card_height + 0.3)
            
            self._add_value_card(slide, card, x, y, card_width, card_height)

    def _add_value_card(self, slide, card: Dict, x: float, y: float, 
                        width: float, height: float) -> None:
        """Add a single value card with icon, title, description, and before/after."""
        icon = card.get('icon', '📊')
        title = card.get('title', '')
        description = card.get('description', '')
        before = card.get('before', '')
        after = card.get('after', '')
        
        # Card background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = self._hex_to_rgb(self.COLORS["light_gray"])
        shape.line.width = Pt(1)
        shape.shadow.inherit = False
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(width), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(32)
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.65), Inches(width - 0.2), Inches(0.4))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_blue"])
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        if description:
            desc_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.05), Inches(width - 0.2), Inches(0.6))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = description
            p.font.size = Pt(10)
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["medium_gray"])
            p.alignment = PP_ALIGN.CENTER
        
        # Before/After if provided
        if before and after:
            ba_y = y + height - 0.6
            
            # Before (red, strikethrough)
            before_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(ba_y), Inches(1.2), Inches(0.35))
            tf = before_box.text_frame
            p = tf.paragraphs[0]
            p.text = before
            p.font.size = Pt(11)
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_red"])
            p.alignment = PP_ALIGN.CENTER
            
            # Arrow
            arrow_box = slide.shapes.add_textbox(Inches(x + 1.4), Inches(ba_y), Inches(0.6), Inches(0.35))
            tf = arrow_box.text_frame
            p = tf.paragraphs[0]
            p.text = "→"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_green"])
            p.alignment = PP_ALIGN.CENTER
            
            # After (green)
            after_box = slide.shapes.add_textbox(Inches(x + 2.0), Inches(ba_y), Inches(1.5), Inches(0.35))
            tf = after_box.text_frame
            p = tf.paragraphs[0]
            p.text = after
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_green"])
            p.alignment = PP_ALIGN.CENTER

    def _populate_before_after_slide(self, slide, config: Dict) -> None:
        """Populate a slide showing before/after transformation."""
        # Remove template placeholders first
        self._remove_placeholder_shapes(slide)
        
        title = config.get('title', '')
        items = config.get('items', [])
        
        # Set title
        self._add_title_textbox(slide, title, 0.5, 0.3, 12.333, size=28)
        
        # Create table-like layout with before/after
        y_start = 1.3
        row_height = 0.7
        
        # Headers
        self._add_ba_header(slide, "Challenge", 0.5, y_start, 4.5, "ms_red")
        self._add_ba_header(slide, "", 5.1, y_start, 1.0, "ms_green")  # Arrow column
        self._add_ba_header(slide, "Solution", 6.2, y_start, 6.5, "ms_green")
        
        for i, item in enumerate(items):
            row_y = y_start + 0.5 + (i * row_height)
            before = item.get('before', '')
            after = item.get('after', '')
            
            # Before text
            self._add_ba_item(slide, before, 0.5, row_y, 4.5, "dark_gray")
            
            # Arrow
            arrow = slide.shapes.add_textbox(Inches(5.1), Inches(row_y), Inches(1.0), Inches(0.5))
            tf = arrow.text_frame
            p = tf.paragraphs[0]
            p.text = "→"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_green"])
            p.alignment = PP_ALIGN.CENTER
            
            # After text
            self._add_ba_item(slide, after, 6.2, row_y, 6.5, "ms_green", bold=True)

    def _add_ba_header(self, slide, text: str, x: float, y: float, 
                       width: float, color: str) -> None:
        """Add a before/after header."""
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.45))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS[color])

    def _add_ba_item(self, slide, text: str, x: float, y: float, 
                     width: float, color: str, bold: bool = False) -> None:
        """Add a before/after item."""
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {text}"
        p.font.size = Pt(14)
        p.font.bold = bold
        p.font.color.rgb = self._hex_to_rgb(self.COLORS[color])

    def _populate_agent_cards_slide(self, slide, config: Dict) -> None:
        """Populate a slide with agent cards (colored boxes like HTML demo)."""
        # Remove template placeholders first
        self._remove_placeholder_shapes(slide)
        
        title = config.get('title', '')
        agents = config.get('agents', [])
        
        # Set title
        self._add_title_textbox(slide, title, 0.5, 0.3, 12.333, size=28)
        
        # Layout agents in a grid
        num_agents = len(agents)
        cols = min(3, num_agents)
        card_width = 3.9
        card_height = 2.0
        gap = 0.2
        
        start_x = (13.333 - (cols * card_width + (cols - 1) * gap)) / 2
        
        for i, agent in enumerate(agents):
            row = i // cols
            col = i % cols
            x = start_x + col * (card_width + gap)
            y = 1.3 + row * (card_height + 0.2)
            
            self._add_agent_card(slide, agent, x, y, card_width, card_height)

    def _add_agent_card(self, slide, agent: Dict, x: float, y: float, 
                        width: float, height: float) -> None:
        """Add a single agent card with gradient-like appearance."""
        name = agent.get('name', '')
        level = agent.get('level', 1)
        description = agent.get('description', '')
        competitors = agent.get('competitors', [])
        
        # Color based on level
        if level == 0:
            bg_color = "#11998e"  # Green for orchestrator
        elif level == 2:
            bg_color = "#667eea"  # Purple for synthesizer
        else:
            bg_color = "#0078d4"  # Blue for Level 1
        
        # Card background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._hex_to_rgb(bg_color)
        shape.line.fill.background()
        
        # Level badge
        badge = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(1.0), Inches(0.25))
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = f"LEVEL {level}"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Agent name
        name_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.4), Inches(width - 0.2), Inches(0.4))
        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Description
        if description:
            desc_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.85), Inches(width - 0.2), Inches(0.6))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = description
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(230, 230, 230)
        
        # Competitors (if any)
        if competitors:
            comp_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + height - 0.4), Inches(width - 0.2), Inches(0.3))
            tf = comp_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = " | ".join(competitors[:4])  # Max 4 competitors
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(200, 200, 200)

    def _populate_metric_boxes_slide(self, slide, config: Dict) -> None:
        """Populate a slide with large metric/stat boxes."""
        # Remove template placeholders first
        self._remove_placeholder_shapes(slide)
        
        title = config.get('title', '')
        metrics = config.get('metrics', [])
        
        # Set title
        self._add_title_textbox(slide, title, 0.5, 0.3, 12.333, size=28)
        
        # Calculate positions
        num_metrics = len(metrics)
        box_width = 3.5
        box_height = 2.5
        gap = 0.4
        
        total_width = num_metrics * box_width + (num_metrics - 1) * gap
        start_x = (13.333 - total_width) / 2
        
        for i, metric in enumerate(metrics):
            x = start_x + i * (box_width + gap)
            self._add_metric_box(slide, metric, x, 2.0, box_width, box_height)

    def _add_metric_box(self, slide, metric: Dict, x: float, y: float,
                        width: float, height: float) -> None:
        """Add a single metric box with large number and label."""
        value = metric.get('value', '')
        label = metric.get('label', '')
        description = metric.get('description', '')
        color = metric.get('color', 'ms_blue')
        
        # Box background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS.get(color, self.COLORS["ms_blue"]))
        shape.line.fill.background()
        
        # Large value
        val_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.3), Inches(width), Inches(1.0))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(y + 1.4), Inches(width), Inches(0.5))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        if description:
            desc_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.9), Inches(width - 0.2), Inches(0.5))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = description
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(220, 220, 220)
            p.alignment = PP_ALIGN.CENTER

    def _populate_process_flow_slide(self, slide, config: Dict) -> None:
        """Populate a slide with a horizontal process flow."""
        # Remove template placeholders first
        self._remove_placeholder_shapes(slide)
        
        title = config.get('title', '')
        steps = config.get('steps', [])
        
        # Set title
        self._add_title_textbox(slide, title, 0.5, 0.3, 12.333, size=28)
        
        num_steps = len(steps)
        if num_steps == 0:
            return
            
        # Calculate positions
        step_width = 2.0
        arrow_width = 0.8
        total_width = num_steps * step_width + (num_steps - 1) * arrow_width
        start_x = (13.333 - total_width) / 2
        
        for i, step in enumerate(steps):
            x = start_x + i * (step_width + arrow_width)
            self._add_process_step(slide, step, x, 2.5, step_width, i + 1)
            
            # Add arrow between steps
            if i < num_steps - 1:
                arrow_x = x + step_width + 0.1
                self._add_flow_arrow(slide, arrow_x, 3.5, arrow_width - 0.2)

    def _add_process_step(self, slide, step: Dict, x: float, y: float,
                          width: float, number: int) -> None:
        """Add a single process step with number circle and description."""
        title = step.get('title', step) if isinstance(step, dict) else step
        description = step.get('description', '') if isinstance(step, dict) else ''
        duration = step.get('duration', '') if isinstance(step, dict) else ''
        
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + width/2 - 0.3), Inches(y), Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["ms_blue"])
        circle.line.fill.background()
        
        # Number text
        num_box = slide.shapes.add_textbox(Inches(x + width/2 - 0.3), Inches(y + 0.08), Inches(0.6), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.8), Inches(width), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_blue"])
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        if description:
            desc_box = slide.shapes.add_textbox(Inches(x), Inches(y + 1.4), Inches(width), Inches(0.8))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = description
            p.font.size = Pt(9)
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["medium_gray"])
            p.alignment = PP_ALIGN.CENTER
        
        # Duration badge
        if duration:
            dur_box = slide.shapes.add_textbox(Inches(x + width/2 - 0.4), Inches(y + 2.2), Inches(0.8), Inches(0.3))
            tf = dur_box.text_frame
            p = tf.paragraphs[0]
            p.text = duration
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_orange"])
            p.alignment = PP_ALIGN.CENTER

    def _add_flow_arrow(self, slide, x: float, y: float, width: float) -> None:
        """Add a flow arrow between process steps."""
        arrow = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.5))
        tf = arrow.text_frame
        p = tf.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_blue"])
        p.alignment = PP_ALIGN.CENTER

    def _populate_bullets(self, text_frame, bullets: List[str]) -> None:
        """Populate a text frame with bullet points."""
        # Clear existing paragraphs except first
        while len(text_frame.paragraphs) > 1:
            p = text_frame.paragraphs[-1]._p
            text_frame._txBody.remove(p)

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.paragraphs[-1]._p.addnext(text_frame.paragraphs[0]._p.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}p', {}))
                p = text_frame.paragraphs[-1]

            p.text = bullet
            self._style_text(p, "body")

    def _style_text(self, paragraph, style: str) -> None:
        """Apply font styling to a paragraph."""
        font_config = self.FONTS.get(style, self.FONTS["body"])
        if paragraph.runs:
            for run in paragraph.runs:
                run.font.name = font_config["name"]
                run.font.size = Pt(font_config["size"])
        else:
            paragraph.font.name = font_config["name"]
            paragraph.font.size = Pt(font_config["size"])

    # ==================== HELPER METHODS FOR CONTENT ====================

    def _add_title_textbox(self, slide, text: str, x: float, y: float, width: float, 
                            size: int = 44, color: str = None) -> None:
        """Add a title text box."""
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(1))
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.FONTS["title"]["name"]
        p.font.size = Pt(size)
        # Use provided color or default to dark_gray
        text_color = color if color else self.COLORS["dark_gray"]
        p.font.color.rgb = self._hex_to_rgb(text_color)

    def _add_subtitle_textbox(self, slide, text: str, x: float, y: float, width: float,
                              color: str = None) -> None:
        """Add a subtitle text box."""
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.6))
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.FONTS["subtitle"]["name"]
        p.font.size = Pt(self.FONTS["subtitle"]["size"])
        # Use provided color or default to medium_gray
        text_color = color if color else self.COLORS["medium_gray"]
        p.font.color.rgb = self._hex_to_rgb(text_color)

    def _add_bullet_textbox(self, slide, bullets: List[str], x: float, y: float, 
                            width: float, height: float, color: str = None) -> None:
        """Add a text box with bullet points."""
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        tf = textbox.text_frame
        tf.word_wrap = True
        # Use provided color or default to dark_gray
        text_color = color if color else self.COLORS["dark_gray"]

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.name = self.FONTS["body"]["name"]
            p.font.size = Pt(self.FONTS["body"]["size"])
            p.font.color.rgb = self._hex_to_rgb(text_color)
            p.space_after = Pt(12)

    def _add_two_column_content(self, slide, left_title: str, right_title: str,
                                left_items: List[str], right_items: List[str]) -> None:
        """Add two-column content with titles and bullet points."""
        # Left column title
        left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.5))
        tf = left_header.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.name = self.FONTS["body"]["name"]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_blue"])
        
        # Right column title
        right_header = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.8), Inches(0.5))
        tf = right_header.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.name = self.FONTS["body"]["name"]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_blue"])
        
        # Left column bullets
        if left_items:
            left_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.8), Inches(4.5))
            tf = left_content.text_frame
            tf.word_wrap = True
            for i, item in enumerate(left_items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.name = self.FONTS["body"]["name"]
                p.font.size = Pt(16)
                p.font.color.rgb = self._hex_to_rgb(self.COLORS["dark_gray"])
                p.space_after = Pt(10)
        
        # Right column bullets
        if right_items:
            right_content = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(5.8), Inches(4.5))
            tf = right_content.text_frame
            tf.word_wrap = True
            for i, item in enumerate(right_items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.name = self.FONTS["body"]["name"]
                p.font.size = Pt(16)
                p.font.color.rgb = self._hex_to_rgb(self.COLORS["dark_gray"])
                p.space_after = Pt(10)

    def _add_comparison_content(self, slide, left_label: str, right_label: str,
                                left_items: List[str], right_items: List[str]) -> None:
        """Add comparison content to a slide."""
        # Left column header
        self._add_column_header(slide, left_label, 0.5, 1.8, 5.5, "ms_red")
        # Right column header
        self._add_column_header(slide, right_label, 7.0, 1.8, 5.5, "ms_blue")

        # Left items
        y_start = 2.5
        for i, item in enumerate(left_items):
            self._add_comparison_item(slide, item, 0.5, y_start + (i * 0.7), 5.5, "ms_red")

        # Right items
        for i, item in enumerate(right_items):
            self._add_comparison_item(slide, item, 7.0, y_start + (i * 0.7), 5.5, "ms_blue")

        # Arrow in the middle
        self._add_arrow(slide, 6.0, 3.5)

    def _add_column_header(self, slide, text: str, x: float, y: float, 
                           width: float, color: str) -> None:
        """Add a column header."""
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.5))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = text.upper()
        p.font.name = self.FONTS["body"]["name"]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS[color])

    def _add_comparison_item(self, slide, text: str, x: float, y: float,
                             width: float, color: str) -> None:
        """Add a comparison item with bullet."""
        # Bullet circle
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y + 0.1),
            Inches(0.15), Inches(0.15)
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS[color])
        bullet.line.fill.background()

        # Text
        textbox = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y), Inches(width - 0.25), Inches(0.6))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.FONTS["body"]["name"]
        p.font.size = Pt(16)
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["dark_gray"])

    def _add_arrow(self, slide, x: float, y: float) -> None:
        """Add an arrow shape."""
        textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1), Inches(0.5))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(36)
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_blue"])
        p.alignment = PP_ALIGN.CENTER

    def _add_quote_box(self, slide, quote: str, author: str) -> None:
        """Add a quote box."""
        # Quote background
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(2),
            Inches(11.333), Inches(3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["ms_blue"])
        box.line.fill.background()

        # Quote text
        textbox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.333), Inches(2))
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'"{quote}"'
        p.font.name = "Segoe UI Light"
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.color.rgb = self._hex_to_rgb(self.COLORS["white"])
        p.alignment = PP_ALIGN.CENTER

        # Author
        if author:
            author_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.333), Inches(0.5))
            tf = author_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {author}"
            p.font.name = self.FONTS["caption"]["name"]
            p.font.size = Pt(16)
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["medium_gray"])
            p.alignment = PP_ALIGN.RIGHT

    def _add_stats_boxes(self, slide, prs, stats: List[Dict]) -> None:
        """Add statistics/metric boxes."""
        num_stats = len(stats)
        if num_stats == 0:
            return

        box_width = min(3.5, 11.0 / num_stats)
        spacing = (12.333 - (box_width * num_stats)) / (num_stats + 1)
        y_start = 2.0

        for i, stat in enumerate(stats):
            x = 0.5 + spacing + (i * (box_width + spacing))
            value = stat.get('value', '')
            label = stat.get('label', '')
            sublabel = stat.get('sublabel', '')

            # Box background
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y_start),
                Inches(box_width), Inches(2.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["white"])
            box.line.color.rgb = self._hex_to_rgb(self.COLORS["light_gray"])
            box.line.width = Pt(1)

            # Value
            value_box = slide.shapes.add_textbox(
                Inches(x), Inches(y_start + 0.4),
                Inches(box_width), Inches(0.8)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.name = "Segoe UI Light"
            p.font.size = Pt(48)
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["ms_blue"])
            p.alignment = PP_ALIGN.CENTER

            # Label
            label_box = slide.shapes.add_textbox(
                Inches(x), Inches(y_start + 1.4),
                Inches(box_width), Inches(0.5)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.name = self.FONTS["body"]["name"]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["dark_gray"])
            p.alignment = PP_ALIGN.CENTER

            # Sublabel
            if sublabel:
                sub_box = slide.shapes.add_textbox(
                    Inches(x), Inches(y_start + 1.9),
                    Inches(box_width), Inches(0.5)
                )
                tf = sub_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = sublabel
                p.font.name = self.FONTS["caption"]["name"]
                p.font.size = Pt(12)
                p.font.color.rgb = self._hex_to_rgb(self.COLORS["medium_gray"])
                p.alignment = PP_ALIGN.CENTER

    def _add_pipeline_boxes(self, slide, prs, steps: List[Dict]) -> None:
        """Add pipeline/process flow boxes."""
        num_steps = len(steps)
        if num_steps == 0:
            return

        # Calculate dimensions
        total_width = 12.333
        step_width = (total_width - 1) / num_steps
        y_start = 2.5

        for i, step in enumerate(steps):
            x = 0.5 + (i * step_width)
            label = step.get('label', f'Step {i+1}')
            description = step.get('description', '')
            number = step.get('number', i + 1)

            # Circle with number
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + (step_width/2) - 0.3), Inches(y_start),
                Inches(0.6), Inches(0.6)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self._hex_to_rgb(self.COLORS["ms_blue"])
            circle.line.fill.background()

            # Number in circle
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.text = str(number)
            p.font.name = self.FONTS["body"]["name"]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["white"])
            p.alignment = PP_ALIGN.CENTER

            # Label below circle
            label_box = slide.shapes.add_textbox(
                Inches(x), Inches(y_start + 0.8),
                Inches(step_width), Inches(0.5)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.name = self.FONTS["body"]["name"]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb(self.COLORS["dark_gray"])
            p.alignment = PP_ALIGN.CENTER

            # Description
            if description:
                desc_box = slide.shapes.add_textbox(
                    Inches(x), Inches(y_start + 1.3),
                    Inches(step_width), Inches(0.5)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = description
                p.font.name = self.FONTS["caption"]["name"]
                p.font.size = Pt(12)
                p.font.color.rgb = self._hex_to_rgb(self.COLORS["medium_gray"])
                p.alignment = PP_ALIGN.CENTER

            # Arrow to next step
            if i < num_steps - 1:
                arrow_x = x + step_width - 0.3
                arrow_box = slide.shapes.add_textbox(
                    Inches(arrow_x), Inches(y_start + 0.1),
                    Inches(0.5), Inches(0.5)
                )
                tf = arrow_box.text_frame
                p = tf.paragraphs[0]
                p.text = "→"
                p.font.size = Pt(24)
                p.font.color.rgb = self._hex_to_rgb(self.COLORS["light_gray"])
                p.alignment = PP_ALIGN.CENTER

    def _save_presentation(self, prs: Presentation, filename: str, kwargs: Dict) -> str:
        """Save the presentation to file."""
        if not filename.endswith('.pptx'):
            filename = f"{filename}.pptx"

        # Get customer name for subfolder organization
        customer = kwargs.get('customer', '')
        
        # Build output directory - create customer subfolder if specified
        base_output_dir = kwargs.get('output_dir', os.path.join(self.base_path, 'docs', 'ppt'))
        if customer:
            # Sanitize customer name for folder (lowercase, replace spaces with underscores)
            customer_folder = customer.lower().replace(' ', '_').replace('-', '_')
            customer_folder = ''.join(c for c in customer_folder if c.isalnum() or c == '_')
            output_dir = os.path.join(base_output_dir, customer_folder)
        else:
            output_dir = base_output_dir
            
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        local_path = os.path.join(output_dir, filename)

        try:
            prs.save(local_path)
        except Exception as e:
            # Fallback to temp directory
            import tempfile
            local_path = os.path.join(tempfile.gettempdir(), filename)
            prs.save(local_path)

        result = {
            "status": "success",
            "filename": filename,
            "local_path": os.path.abspath(local_path),
            "customer": customer if customer else "none",
            "customer_folder": output_dir,
            "template_used": kwargs.get('template', 'BaseTemplateBlue')
        }

        return json.dumps(result, indent=2)
