"""Microsoft Deck Studio — generate a polished, Microsoft-Fluent-styled PowerPoint (.pptx)
from a plain JSON spec of slides. Give it a deck spec; it renders a branded 16:9 deck
(Segoe UI, Fluent palette, Microsoft 4-square mark, brand bar + page footer) and returns
the saved file path. python-pptx is imported lazily (auto-installed if missing).

SPEC SHAPE (pass as the `spec` param, JSON string or object):
{
  "deck": { "wordmark":"RAPP · Rapid Agent Prototype Platform",
             "footLeft":"Microsoft · MCAPS", "pageLabel":"Business Overview" },
  "slides": [ { "layout":"<name>", ... }, ... ]
}

LAYOUTS (each slide is one object; {h}...{/h} in a headline highlights that span in the accent color):
  title      : title, titleSize?, expand?, tag?, one?            (navy hero cover)
  hero       : kicker, hook, subhead?, proofChips[]?, closingLine?   (dark hook slide)
  statement  : kicker?, headline({h}), sub? OR points[]           (big statement + supports)
  bullets    : kicker?, headline, points[]  ("Lead||rest" bolds the lead)
  steps      : kicker?, headline, steps[{n?,t,d}] (<=3), footnote?  (numbered cards)
  columns    : kicker?, headline, columns[{tag,title,points[],alt?}] (<=2), footnote?
  content2col: kicker?, headline, intro?, accent?, columns[{title,points[],accent?}], highlight?
  feature    : kicker, headline({h}), points[], accent?, side?(left|right),
               visual?(bubble|number|file), hook?, hookSub?, bubbleText?
  ecosystem  : kicker?, headline({h}), subhead?, icons[{file,label}]  (file = local PNG path)
  flow       : kicker?, headline, subhead?, stages[{name,sub,owner?,accent?}], footer?
  processflow: kicker?, headline, steps[{title,desc}], highlights[<=2], footer?
  quote      : quote, by?, role?
  cta        : kicker?, headline, points[]?, doorLabel?, link?
Colors accept hex like "#0078D4" (blue), "#8661C5" (purple), "#D83B01" (orange), "#107C10" (green).
"""

__manifest__ = {
    "schema": "rapp-agent/1.0",
    "name": "@kody-w/microsoft_deck_studio_agent",
    "version": "1.0.1",
    "display_name": "Microsoft Deck Studio",
    "description": "Generate a polished, Microsoft-Fluent-styled PowerPoint (.pptx) deck from a JSON spec of slides and layouts.",
    "author": "Kody Wildfeuer",
    "tags": ["powerpoint", "pptx", "deck", "slides", "presentation", "microsoft", "fluent", "generator", "design"],
    "category": "productivity",
    "quality_tier": "community",
    "requires_env": [],
    "dependencies": ["@rapp/basic_agent"],
}

import os, sys, json, math, re

try:
    from agents.basic_agent import BasicAgent            # RAPP runtime
except ImportError:
    try:
        from basic_agent import BasicAgent
    except ImportError:                                   # standalone fallback
        class BasicAgent:
            def __init__(self, name=None, metadata=None):
                self.name = name or "BasicAgent"; self.metadata = metadata or {}
            def perform(self, **kwargs): return "Not implemented."
            def system_context(self): return None

# ---- brand strings (set per-render from spec["deck"]) ----
FONT = "Segoe UI"; FONT_SB = "Segoe UI Semibold"
WORDMARK = "RAPP · Rapid Agent Prototype Platform"
FOOT_LEFT = "Microsoft · MCAPS"; FOOT_RIGHT = "Business Overview · v1"; PAGE_LABEL = "Business Overview"
_ENGINE_READY = False


def _ensure_engine():
    """Import python-pptx (auto-install if missing) and bind pptx symbols + palette as globals."""
    global _ENGINE_READY, Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, qn
    global BLUE, BLUE_DK, NAVY, NAVY_DEEP, CYAN, INK, INK_SOFT, MUTED, LINE, SURFACE, WHITE, PAPER, SQ, LIGHTBLUE, PALEBLUE
    if _ENGINE_READY:
        return True, ""

    def _imp():
        global Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, qn
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn
    try:
        _imp()
    except ImportError:
        try:
            import subprocess
            subprocess.check_call([sys.executable, "-m", "pip", "install", "--quiet", "python-pptx"])
            _imp()
        except Exception as e:
            return False, ("[MicrosoftDeckStudio] needs python-pptx and could not auto-install it. "
                           "Run:  pip install python-pptx   (%s)" % e)
    BLUE = RGBColor(0x00, 0x78, 0xD4); BLUE_DK = RGBColor(0x00, 0x5A, 0x9E); NAVY = RGBColor(0x10, 0x3A, 0x6B)
    NAVY_DEEP = RGBColor(0x0B, 0x2A, 0x4A); CYAN = RGBColor(0x50, 0xE6, 0xFF); INK = RGBColor(0x1B, 0x1A, 0x19)
    INK_SOFT = RGBColor(0x3B, 0x3A, 0x39); MUTED = RGBColor(0x60, 0x5E, 0x5C); LINE = RGBColor(0xE1, 0xDF, 0xDD)
    SURFACE = RGBColor(0xF3, 0xF2, 0xF1); WHITE = RGBColor(0xFF, 0xFF, 0xFF); PAPER = RGBColor(0xFF, 0xFF, 0xFF)
    SQ = [RGBColor(0xF2, 0x50, 0x22), RGBColor(0x7F, 0xBA, 0x00), RGBColor(0x00, 0xA4, 0xEF), RGBColor(0xFF, 0xB9, 0x00)]
    LIGHTBLUE = RGBColor(0xBF, 0xE6, 0xFF); PALEBLUE = RGBColor(0xDC, 0xEE, 0xFF)
    _ENGINE_READY = True
    return True, ""


# ---------- color helpers ----------
def _rgb(h):
    h = h.lstrip('#'); return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
def C(h):
    r, g, b = _rgb(h); return RGBColor(r, g, b)
def lighten(h, f):
    r, g, b = _rgb(h); return RGBColor(int(r + (255 - r) * f), int(g + (255 - g) * f), int(b + (255 - b) * f))
def darken(h, f):
    r, g, b = _rgb(h); return RGBColor(int(r * (1 - f)), int(g * (1 - f)), int(b * (1 - f)))

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, line_w=None, shape=None, shadow=False):
    if shape is None: shape = MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    if shadow: _soft_shadow(sp)
    return sp

def _soft_shadow(sp):
    spPr = sp._element.spPr
    for ex in spPr.findall(qn('a:effectLst')): spPr.remove(ex)
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'), {'blurRad': '90000', 'dist': '40000', 'dir': '5400000', 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = clr.makeelement(qn('a:alpha'), {'val': '22000'})
    clr.append(alpha); sh.append(clr); el.append(sh); spPr.append(el)

def grad(sp, c1, c2, angle=45):
    sp.fill.gradient(); stops = sp.fill.gradient_stops
    stops[0].position = 0.0; stops[0].color.rgb = c1
    stops[1].position = 1.0; stops[1].color.rgb = c2
    try: sp.fill.gradient_angle = angle
    except Exception: pass

def txt(s, x, y, w, h, runs, align=None, anchor=None, space_after=None, line_spacing=None, wrap=True):
    if align is None: align = PP_ALIGN.LEFT
    if anchor is None: anchor = MSO_ANCHOR.TOP
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        al = para.get("align"); p.alignment = al if al is not None else align
        if line_spacing: p.line_spacing = line_spacing
        if para.get("line_spacing"): p.line_spacing = para["line_spacing"]
        if space_after is not None: p.space_after = Pt(space_after)
        if para.get("space_before") is not None: p.space_before = Pt(para["space_before"])
        if para.get("space_after") is not None: p.space_after = Pt(para["space_after"])
        for rdef in para["runs"]:
            r = p.add_run(); r.text = rdef["t"]; f = r.font
            f.size = Pt(rdef.get("sz", 18)); f.name = rdef.get("font") or FONT
            f.bold = rdef.get("b", False); f.color.rgb = rdef.get("c") if rdef.get("c") is not None else INK
            if rdef.get("spacing") is not None: _letter_spacing(r, rdef["spacing"])
    return tb

def _letter_spacing(run, pts):
    run._r.get_or_add_rPr().set('spc', str(int(pts * 100)))

def R(t, sz=18, c=None, b=False, font=None, spacing=None):
    return {"t": t, "sz": sz, "c": c, "b": b, "font": font, "spacing": spacing}
def P(runs, align=None, sa=None, sb=None, ls=None):
    d = {"runs": runs, "align": align}
    if sa is not None: d["space_after"] = sa
    if sb is not None: d["space_before"] = sb
    if ls is not None: d["line_spacing"] = ls
    return d

def squares(s, x, y, size, gap=None):
    g = Emu(int(size * 0.14)); pos = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for i, (cx, cy) in enumerate(pos):
        rect(s, x + cx * (size + g), y + cy * (size + g), size, size, fill=SQ[i], shape=MSO_SHAPE.RECTANGLE)

def brandbar(s, dark=False):
    x = Inches(0.55); y = Inches(0.42)
    squares(s, x, y, Inches(0.13))
    name, _, desc = WORDMARK.partition(" · ")
    txt(s, x + Inches(0.42), y - Inches(0.04), Inches(6.5), Inches(0.4),
        [P([R(name + " ", 14, INK if not dark else WHITE, True),
            R(("· " + desc) if desc else "", 12, MUTED if not dark else LIGHTBLUE)])], anchor=MSO_ANCHOR.MIDDLE)

def pagenum(s, i, dark=False):
    txt(s, SW - Inches(3.1), SH - Inches(0.55), Inches(2.6), Inches(0.35),
        [P([R("%s · %02d" % (PAGE_LABEL, i), 10.5, MUTED if not dark else LIGHTBLUE, spacing=0.6)], align=PP_ALIGN.RIGHT)],
        anchor=MSO_ANCHOR.MIDDLE)

def kicker(s, x, y, text, dark=False):
    rect(s, x, y + Inches(0.10), Inches(0.30), Inches(0.045), fill=BLUE if not dark else CYAN)
    txt(s, x + Inches(0.42), y, Inches(9), Inches(0.35),
        [P([R(text.upper(), 13, BLUE if not dark else CYAN, True, spacing=1.6)])], anchor=MSO_ANCHOR.MIDDLE)

def bg(s, color): rect(s, 0, 0, SW, SH, fill=color)
def bg_grad(s, c1, c2, angle=45):
    sp = rect(s, 0, 0, SW, SH, fill=c1); grad(sp, c1, c2, angle); return sp

def check_list(s, x, y, w, items, gap=None, size=20, dark=False, line_h=None, chk_color=None, chk_bg=None):
    if line_h is None: line_h = Inches(0.72)
    cc = chk_color or (WHITE if dark else BLUE)
    cbg = chk_bg or (RGBColor(0x1C, 0x4E, 0x82) if dark else RGBColor(0xE5, 0xF1, 0xFB))
    cy = y
    for it in items:
        rect(s, x, cy + Inches(0.03), Inches(0.34), Inches(0.34), fill=cbg, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x, cy + Inches(0.03), Inches(0.34), Inches(0.34),
            [P([R("✓", 14, cc, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        lead, rest = (it.split("||", 1) + [""])[:2] if "||" in it else ("", it)
        runs = []
        if lead: runs.append(R(lead + " ", size, (INK if not dark else WHITE), True))
        runs.append(R(rest, size, (INK_SOFT if not dark else PALEBLUE)))
        txt(s, x + Inches(0.52), cy - Inches(0.02), w - Inches(0.52), line_h, [P(runs, ls=1.08)], anchor=MSO_ANCHOR.TOP)
        cy = cy + line_h
    return cy

# ============================ layout renderers ============================
def render_title(s, d, i):
    bg_grad(s, NAVY_DEEP, RGBColor(0x0E, 0x5A, 0xA7), angle=60)
    rect(s, 0, SH - Inches(0.18), SW, Inches(0.18), fill=BLUE)
    x = Inches(1.0); squares(s, x, Inches(1.1), Inches(0.32))
    tsize = d.get("titleSize", 100); title = d.get("title", "RAPP")
    lh = tsize * 1.08 / 72.0; cpl = max(1, int(11.2 / ((tsize * 0.52) / 72.0)))
    tlines = max(1, math.ceil(len(title) / cpl))
    txt(s, x, Inches(1.62), Inches(11.5), Inches(tlines * lh + 0.3),
        [P([R(title, tsize, WHITE, True, font=FONT_SB)], ls=1.06)])
    cur = 1.62 + tlines * lh + 0.34
    txt(s, x + Inches(0.03), Inches(cur), Inches(11), Inches(0.5),
        [P([R(d.get("expand", "Rapid Agent Prototype Platform"), 20, LIGHTBLUE, True, spacing=0.8)])])
    cur += 0.6
    if d.get("tag"):
        txt(s, x, Inches(cur), Inches(10.8), Inches(1.0), [P([R(d["tag"], 24, RGBColor(0xEA, 0xF6, 0xFF), True)], ls=1.2)]); cur += 0.9
    if d.get("one"):
        txt(s, x, Inches(cur), Inches(10.6), Inches(1.3), [P([R(d["one"], 15.5, RGBColor(0x9F, 0xCB, 0xEE))], ls=1.3)])
    txt(s, x, SH - Inches(0.72), Inches(6.5), Inches(0.4), [P([R(FOOT_LEFT, 13, RGBColor(0x9F, 0xC6, 0xE8), True, spacing=0.8)])])
    txt(s, SW - Inches(4.0), SH - Inches(0.72), Inches(3.0), Inches(0.4),
        [P([R(FOOT_RIGHT, 13, RGBColor(0x9F, 0xC6, 0xE8), spacing=0.8)], align=PP_ALIGN.RIGHT)])

def render_statement(s, d, i):
    bg(s, PAPER); rect(s, 0, 0, Inches(0.16), SH, fill=BLUE); brandbar(s); pagenum(s, i)
    x = Inches(1.0)
    if d.get("kicker"): kicker(s, x, Inches(2.0), d["kicker"])
    head = d.get("headline", ""); runs = []
    for part in re.split(r"(\{h\}.*?\{/h\})", head):
        if part.startswith("{h}"): runs.append(R(part[3:-4], 46, BLUE, True, font=FONT_SB))
        elif part: runs.append(R(part, 46, INK, True, font=FONT_SB))
    txt(s, x, Inches(2.45), Inches(11.4), Inches(1.9), [P(runs, ls=1.06)])
    plain = re.sub(r"\{/?h\}", "", head); hlines = max(1, math.ceil(len(plain) / 34))
    body_top = 2.45 + hlines * 0.70 + 0.42
    if d.get("sub"):
        txt(s, x, Inches(body_top), Inches(10.6), Inches(1.2), [P([R(d["sub"], 22, MUTED)], ls=1.3)])
    if d.get("points"):
        cy = Inches(body_top)
        for pt in d["points"]:
            nl = max(1, math.ceil(len(pt) / max(1, int(10.7 / ((18.5 * 0.52) / 72.0)))))
            rect(s, x, cy + Inches(0.12), Inches(0.13), Inches(0.13), fill=BLUE, shape=MSO_SHAPE.OVAL)
            txt(s, x + Inches(0.34), cy, Inches(10.7), Inches(0.32 * nl + 0.1), [P([R(pt, 18.5, INK_SOFT)], ls=1.15)], anchor=MSO_ANCHOR.TOP)
            cy = cy + Inches(0.30 * nl + 0.32)

def render_bullets(s, d, i):
    bg(s, PAPER); brandbar(s); pagenum(s, i); x = Inches(1.0)
    if d.get("kicker"): kicker(s, x, Inches(1.15), d["kicker"])
    txt(s, x, Inches(1.65), Inches(11.2), Inches(1.4), [P([R(d.get("headline", ""), 38, INK, True, font=FONT_SB)], ls=1.05)])
    pts = d.get("points", []); n = len(pts)
    if n <= 3:   size, line_h, y0 = 23, Inches(1.02), Inches(3.45)
    elif n == 4: size, line_h, y0 = 21, Inches(0.88), Inches(3.15)
    else:        size, line_h, y0 = 19.5, Inches(0.76), Inches(3.0)
    check_list(s, x, y0, Inches(11.0), pts, size=size, line_h=line_h)

def render_steps(s, d, i):
    bg(s, PAPER); brandbar(s); pagenum(s, i); x = Inches(1.0)
    if d.get("kicker"): kicker(s, x, Inches(1.15), d["kicker"])
    txt(s, x, Inches(1.65), Inches(11.2), Inches(1.2), [P([R(d.get("headline", ""), 36, INK, True, font=FONT_SB)], ls=1.05)])
    steps = d.get("steps", []); gap = Inches(0.4); total_w = Inches(11.33); cw = (total_w - gap * 2) / 3
    y = Inches(3.1); ch = Inches(3.3)
    for idx, st in enumerate(steps[:3]):
        cx = x + idx * (cw + gap)
        rect(s, cx, y, cw, ch, fill=WHITE, line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        rect(s, cx, y + Inches(0.2), Inches(0.09), ch - Inches(0.4), fill=BLUE)
        rect(s, cx + Inches(0.42), y + Inches(0.4), Inches(0.7), Inches(0.7), fill=BLUE, shape=MSO_SHAPE.OVAL)
        txt(s, cx + Inches(0.42), y + Inches(0.4), Inches(0.7), Inches(0.7),
            [P([R(st.get("n", str(idx + 1)), 24, WHITE, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, cx + Inches(0.42), y + Inches(1.28), cw - Inches(0.84), Inches(0.82),
            [P([R(st.get("t", ""), 21, INK, True, font=FONT_SB)], ls=1.05)], anchor=MSO_ANCHOR.TOP)
        txt(s, cx + Inches(0.42), y + Inches(2.18), cw - Inches(0.84), Inches(0.9),
            [P([R(st.get("d", ""), 15.5, MUTED)], ls=1.18)], anchor=MSO_ANCHOR.TOP)
    if d.get("footnote"):
        fy = y + ch + Inches(0.32); rect(s, x, fy + Inches(0.04), Inches(0.30), Inches(0.045), fill=BLUE)
        txt(s, x + Inches(0.42), fy - Inches(0.08), Inches(8.3), Inches(0.6), [P([R(d["footnote"], 17, INK_SOFT, True)], ls=1.15)])

def render_columns(s, d, i):
    bg(s, PAPER); brandbar(s); pagenum(s, i); x = Inches(1.0)
    if d.get("kicker"): kicker(s, x, Inches(1.15), d["kicker"])
    txt(s, x, Inches(1.65), Inches(11.2), Inches(1.2), [P([R(d.get("headline", ""), 36, INK, True, font=FONT_SB)], ls=1.05)])
    cols = d.get("columns", []); gap = Inches(0.45); total = Inches(11.33); cw = (total - gap) / 2
    y = Inches(2.9); ch = Inches(3.55); pad = Inches(0.5); inner_w_in = (cw - pad * 2) / 914400.0
    for idx, c in enumerate(cols[:2]):
        cx = x + idx * (cw + gap); alt = c.get("alt", False)
        card = rect(s, cx, y, cw, ch, fill=(NAVY if alt else WHITE), line=(None if alt else LINE), line_w=Pt(1),
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        if alt: grad(card, NAVY, RGBColor(0x0E, 0x5A, 0xA7), 60)
        txt(s, cx + pad, y + Inches(0.42), cw - pad * 2, Inches(0.35),
            [P([R(c.get("tag", "").upper(), 12.5, (CYAN if alt else BLUE), True, spacing=1.4)])])
        txt(s, cx + pad, y + Inches(0.82), cw - pad * 2, Inches(0.7),
            [P([R(c.get("title", ""), 25, (WHITE if alt else INK), True, font=FONT_SB)], ls=1.05)])
        cy = y + Inches(1.65)
        for pt in c.get("points", []):
            nl = max(1, math.ceil(len(pt) / max(1, int(inner_w_in / ((17.5 * 0.52) / 72.0)))))
            rect(s, cx + pad, cy + Inches(0.10), Inches(0.13), Inches(0.13), fill=(CYAN if alt else BLUE), shape=MSO_SHAPE.OVAL)
            txt(s, cx + pad + Inches(0.34), cy - Inches(0.02), cw - pad * 2 - Inches(0.34), Inches(0.34 * nl + 0.1),
                [P([R(pt, 17.5, (PALEBLUE if alt else INK_SOFT))], ls=1.12)], anchor=MSO_ANCHOR.TOP)
            cy = cy + Inches(0.30 * nl + 0.30)
    if d.get("footnote"):
        fy = y + ch + Inches(0.28); rect(s, x, fy + Inches(0.04), Inches(0.30), Inches(0.045), fill=BLUE)
        txt(s, x + Inches(0.42), fy - Inches(0.06), Inches(8.6), Inches(0.5), [P([R(d["footnote"], 16, INK_SOFT, True)], ls=1.1)])

def render_quote(s, d, i):
    bg_grad(s, RGBColor(0xF3, 0xF8, 0xFE), WHITE, 60); brandbar(s); pagenum(s, i); x = Inches(1.35)
    bar = rect(s, Inches(1.0), Inches(2.05), Inches(0.13), Inches(3.0), fill=BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    grad(bar, BLUE, CYAN, 90); squares(s, x, Inches(2.0), Inches(0.13))
    txt(s, x, Inches(2.55), Inches(10.6), Inches(2.6), [P([R(d.get("quote", ""), 38, INK, True, font=FONT_SB)], ls=1.16)], anchor=MSO_ANCHOR.TOP)
    by = d.get("by", ""); role = d.get("role", "")
    rect(s, x, Inches(5.6), Inches(0.30), Inches(0.045), fill=BLUE)
    txt(s, x + Inches(0.42), Inches(5.4), Inches(11), Inches(0.6),
        [P([R(by, 19, BLUE_DK, True), R(("   ·   " + role) if role else "", 17, MUTED)])], anchor=MSO_ANCHOR.MIDDLE)

def render_cta(s, d, i):
    bg_grad(s, NAVY_DEEP, RGBColor(0x0E, 0x5A, 0xA7), 60); rect(s, 0, SH - Inches(0.18), SW, Inches(0.18), fill=BLUE)
    squares(s, Inches(0.55), Inches(0.5), Inches(0.12)); x = Inches(1.0)
    if d.get("kicker"):
        txt(s, x, Inches(1.5), Inches(9), Inches(0.4), [P([R(d["kicker"].upper(), 13, CYAN, True, spacing=1.6)])])
    head = d.get("headline", "")
    txt(s, x, Inches(1.95), Inches(11.4), Inches(1.5), [P([R(head, 48, WHITE, True, font=FONT_SB)], ls=1.05)])
    hlines = max(1, math.ceil(len(head) / 32)); pts_top = 1.95 + hlines * 0.72 + 0.34
    if d.get("points"):
        check_list(s, x, Inches(pts_top), Inches(11.0), d["points"], size=19, dark=True, line_h=Inches(0.66))
    if d.get("link"):
        by = Inches(5.95)
        rect(s, x, by, Inches(7.4), Inches(0.95), fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        rect(s, x + Inches(0.28), by + Inches(0.235), Inches(0.48), Inches(0.48), fill=BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x + Inches(0.28), by + Inches(0.20), Inches(0.48), Inches(0.52),
            [P([R("↓", 22, WHITE, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.95), by, Inches(6.3), Inches(0.95),
            [P([R((d.get("doorLabel", "Start here") + "   "), 19, BLUE_DK, True), R(d.get("link", ""), 18, BLUE, font="Consolas")])], anchor=MSO_ANCHOR.MIDDLE)

def _feature_visual(s, d, px, py, pw, ph, ah):
    accent = C(ah); v = d.get("visual", "number")
    if v == "number":
        txt(s, px, py + Inches(1.25), pw, Inches(2.0), [P([R(d.get("hook", "~30"), 132, WHITE, True, font=FONT_SB)], align=PP_ALIGN.CENTER)])
        if d.get("hookSub"):
            txt(s, px, py + Inches(3.5), pw, Inches(0.7), [P([R(d["hookSub"], 22, WHITE, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    elif v == "bubble":
        bw = pw - Inches(1.0); bx = px + Inches(0.5); byy = py + Inches(1.1); bh = Inches(2.35)
        cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT, bx, byy, bw, bh)
        cb.fill.solid(); cb.fill.fore_color.rgb = WHITE; cb.line.fill.background(); cb.shadow.inherit = False; _soft_shadow(cb)
        tf = cb.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3); tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.35)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = d.get("bubbleText", "“What's my biggest risk today?”")
        r.font.size = Pt(21); r.font.name = FONT_SB; r.font.bold = True; r.font.color.rgb = darken(ah, 0.15)
        if d.get("hookSub"):
            txt(s, px, py + ph - Inches(1.15), pw, Inches(0.8), [P([R(d["hookSub"], 20, WHITE, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    elif v == "file":
        cwd = Inches(2.3); chh = Inches(2.85); cxm = px + pw / 2; fy = py + Inches(1.0)
        rect(s, cxm - cwd / 2 + Inches(0.42), fy + Inches(0.4), cwd, chh, fill=lighten(ah, 0.72), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        fx = cxm - cwd / 2 - Inches(0.18)
        rect(s, fx, fy, cwd, chh, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        rect(s, fx, fy, cwd, Inches(0.6), fill=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, fx, fy + Inches(0.3), cwd, Inches(0.3), fill=accent)
        txt(s, fx + Inches(0.22), fy, cwd - Inches(0.4), Inches(0.6), [P([R(d.get("fileName", "solution.rapp"), 14, WHITE, True, font=FONT_SB)])], anchor=MSO_ANCHOR.MIDDLE)
        ly = fy + Inches(0.92)
        for wln in [1.7, 1.85, 1.45, 1.75]:
            rect(s, fx + Inches(0.24), ly, Inches(wln), Inches(0.12), fill=lighten(ah, 0.55), shape=MSO_SHAPE.ROUNDED_RECTANGLE); ly += Inches(0.42)
        if d.get("hookSub"):
            txt(s, px, py + ph - Inches(1.0), pw, Inches(0.7), [P([R(d["hookSub"], 20, WHITE, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)

def render_feature(s, d, i):
    bg(s, PAPER); ah = d.get("accent", "#0078D4"); accent = C(ah); side = d.get("side", "right")
    pw = Inches(4.7); ph = Inches(5.3); ptop = Inches(1.1)
    if side == "left": px = Inches(0.5); cx = Inches(5.7); cw = Inches(6.9)
    else: px = SW - Inches(5.2); cx = Inches(1.0); cw = Inches(6.6)
    panel = rect(s, px, ptop, pw, ph, fill=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True); grad(panel, accent, darken(ah, 0.30), 60)
    _feature_visual(s, d, px, ptop, pw, ph, ah); brandbar(s); pagenum(s, i)
    ky = Inches(1.5); rect(s, cx, ky + Inches(0.10), Inches(0.30), Inches(0.045), fill=accent)
    txt(s, cx + Inches(0.42), ky, Inches(6), Inches(0.35), [P([R(d.get("kicker", "").upper(), 13, accent, True, spacing=1.6)])], anchor=MSO_ANCHOR.MIDDLE)
    head = d.get("headline", ""); runs = []
    for part in re.split(r"(\{h\}.*?\{/h\})", head):
        if part.startswith("{h}"): runs.append(R(part[3:-4], 40, accent, True, font=FONT_SB))
        elif part: runs.append(R(part, 40, INK, True, font=FONT_SB))
    txt(s, cx, Inches(2.05), cw, Inches(1.7), [P(runs, ls=1.05)])
    plain = re.sub(r"\{/?h\}", "", head); cwin = cw / 914400.0
    hlines = max(1, math.ceil(len(plain) / max(1, int(cwin / ((40 * 0.52) / 72.0)))))
    body_top = 2.05 + hlines * 0.62 + 0.44
    check_list(s, cx, Inches(body_top), cw, d.get("points", []), size=19, line_h=Inches(0.82), chk_color=accent, chk_bg=lighten(ah, 0.86))

def render_hero(s, d, i):
    bg_grad(s, NAVY_DEEP, RGBColor(0x0E, 0x3E, 0x72), angle=55); rect(s, 0, SH - Inches(0.18), SW, Inches(0.18), fill=BLUE)
    brandbar(s, dark=True); pagenum(s, i, dark=True); x = Inches(1.0)
    txt(s, x, Inches(1.5), Inches(10), Inches(0.4), [P([R(d.get("kicker", "THE BIG IDEA").upper(), 14, CYAN, True, spacing=2.2)])])
    hook = d.get("hook", ""); hs = 46
    while hs > 30:
        cpl = max(1, int(11.3 / ((hs * 0.52) / 72.0)))
        if math.ceil(len(hook) / cpl) <= 2: break
        hs -= 2
    cpl = max(1, int(11.3 / ((hs * 0.52) / 72.0))); hlines = max(1, math.ceil(len(hook) / cpl))
    txt(s, x, Inches(2.12), Inches(11.5), Inches(hlines * 0.85 + 0.3), [P([R(hook, hs, WHITE, True, font=FONT_SB)], ls=1.06)])
    cur = 2.12 + hlines * (hs * 1.06 / 72.0) + 0.42
    if d.get("subhead"):
        txt(s, x, Inches(cur), Inches(10.9), Inches(1.0), [P([R(d["subhead"], 20, RGBColor(0xBF, 0xE6, 0xFF))], ls=1.3)])
        cur += 0.42 * max(1, math.ceil(len(d["subhead"]) / 95)) + 0.35
    chip_cols = [BLUE, C("#107C10"), C("#D83B01"), C("#8661C5")]; cxp = Inches(1.0); cyp = Inches(cur + 0.05)
    for idx, ch in enumerate(d.get("proofChips", [])):
        w = Inches(0.55 + 0.125 * len(ch))
        rect(s, cxp, cyp, w, Inches(0.56), fill=chip_cols[idx % len(chip_cols)], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, cxp, cyp, w, Inches(0.56), [P([R(ch, 15, WHITE, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE); cxp = cxp + w + Inches(0.22)
    if d.get("closingLine"):
        txt(s, x, SH - Inches(1.28), Inches(11), Inches(0.7), [P([R(d["closingLine"], 18, CYAN, True)], ls=1.2)])

def render_ecosystem(s, d, i):
    bg(s, PAPER); brandbar(s); pagenum(s, i); x = Inches(1.0)
    if d.get("kicker"): kicker(s, x, Inches(1.05), d["kicker"])
    head = d.get("headline", ""); runs = []
    for part in re.split(r"(\{h\}.*?\{/h\})", head):
        if part.startswith("{h}"): runs.append(R(part[3:-4], 38, BLUE, True, font=FONT_SB))
        elif part: runs.append(R(part, 38, INK, True, font=FONT_SB))
    txt(s, x, Inches(1.5), Inches(11.5), Inches(1.3), [P(runs, ls=1.05)])
    plain = re.sub(r"\{/?h\}", "", head); hlines = max(1, math.ceil(len(plain) / max(1, int(11.5 / ((38 * 0.52) / 72.0)))))
    cur = 1.5 + hlines * 0.58 + 0.22; pw = Inches(4.9)
    rect(s, x, Inches(cur), pw, Inches(0.5), fill=C("#0B2A4A"), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, Inches(cur), pw, Inches(0.5), [P([R(d.get("badge", "●  Runs LOCAL · one file · no cloud to babysit"), 13, WHITE, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    cur += 0.66
    if d.get("subhead"):
        subl = max(1, math.ceil(len(d["subhead"]) / 88))
        txt(s, x, Inches(cur), Inches(11.2), Inches(0.42 * subl + 0.2), [P([R(d["subhead"], 18, MUTED)], ls=1.25)]); cur += 0.42 * subl + 0.26
    txt(s, x, Inches(cur), Inches(11), Inches(0.35), [P([R(d.get("rail", "TRANSLATES INTO"), 13, BLUE, True, spacing=2.2)])])
    icons = d.get("icons", []); n = max(1, len(icons)); band_x = 1.0; band_w = 11.33; colw = band_w / n; iy = cur + 0.48
    for idx, ic in enumerate(icons):
        cxc = band_x + colw * idx + colw / 2
        try:
            pic = s.shapes.add_picture(ic["file"], Inches(cxc - 0.4), Inches(iy), height=Inches(0.8)); pic.left = int(Inches(cxc) - pic.width / 2)
        except Exception:
            rect(s, Inches(cxc - 0.4), Inches(iy), Inches(0.8), Inches(0.8), fill=lighten("#0078D4", 0.85), line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, Inches(band_x + colw * idx), Inches(iy + 0.92), Inches(colw), Inches(0.5), [P([R(ic.get("label", ""), 11, INK_SOFT, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.TOP)

def render_flow(s, d, i):
    bg(s, PAPER); brandbar(s); pagenum(s, i); x = Inches(1.0)
    if d.get("kicker"): kicker(s, x, Inches(1.0), d["kicker"])
    txt(s, x, Inches(1.42), Inches(11.4), Inches(0.9), [P([R(d.get("headline", ""), 40, INK, True, font=FONT_SB)])])
    if d.get("subhead"): txt(s, x, Inches(2.4), Inches(11), Inches(0.6), [P([R(d["subhead"], 19, MUTED)], ls=1.2)])
    stages = d.get("stages", []); n = max(1, len(stages)); band_w = 11.33; arrow_w = 0.34; chip_w = (band_w - (n - 1) * arrow_w) / n
    cy = 3.35; chh = 1.28; cxp = 1.0
    for idx, st in enumerate(stages):
        ah = st.get("accent"); acc = C(ah) if ah else None; fill = lighten(ah, 0.86) if ah else RGBColor(0xF3, 0xF2, 0xF1)
        rect(s, Inches(cxp), Inches(cy), Inches(chip_w), Inches(chh), fill=fill, line=(acc if acc else LINE), line_w=Pt(1.75 if acc else 1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, Inches(cxp), Inches(cy + 0.2), Inches(chip_w), Inches(0.4), [P([R(st["name"].upper(), 13.5, (acc if acc else INK), True, font=FONT_SB)], align=PP_ALIGN.CENTER)])
        txt(s, Inches(cxp + 0.1), Inches(cy + 0.62), Inches(chip_w - 0.2), Inches(0.55), [P([R(st.get("sub", ""), 10.5, MUTED)], align=PP_ALIGN.CENTER, ls=1.05)])
        if st.get("owner"):
            txt(s, Inches(cxp - 0.1), Inches(cy + chh + 0.1), Inches(chip_w + 0.2), Inches(0.35), [P([R("▲ " + st["owner"], 12, acc, True)], align=PP_ALIGN.CENTER)])
        cxp += chip_w
        if idx < n - 1:
            txt(s, Inches(cxp), Inches(cy + 0.32), Inches(arrow_w), Inches(0.5), [P([R("→", 19, MUTED, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE); cxp += arrow_w
    if d.get("footer"):
        fby = 5.5; rect(s, x, Inches(fby), Inches(11.33), Inches(1.2), fill=RGBColor(0xF3, 0xF8, 0xFE), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, Inches(fby), Inches(0.09), Inches(1.2), fill=BLUE)
        txt(s, x + Inches(0.34), Inches(fby), Inches(10.7), Inches(1.2), [P([R(d["footer"], 14.5, INK_SOFT)], ls=1.25)], anchor=MSO_ANCHOR.MIDDLE)

def render_content2col(s, d, i):
    bg(s, PAPER); brandbar(s); pagenum(s, i); ah = d.get("accent", "#0078D4"); accent = C(ah); x = Inches(1.0)
    if d.get("kicker"): kicker(s, x, Inches(1.0), d["kicker"])
    head = d.get("headline", "")
    txt(s, x, Inches(1.42), Inches(11.4), Inches(1.0), [P([R(head, 29, INK, True, font=FONT_SB)], ls=1.03)])
    hl = max(1, math.ceil(len(head) / max(1, int(11.4 / ((29 * 0.52) / 72.0))))); hy = 1.42 + hl * 0.45 + 0.14
    if d.get("intro"):
        il = max(1, math.ceil(len(d["intro"]) / 108))
        txt(s, x, Inches(hy), Inches(11.2), Inches(0.36 * il + 0.15), [P([R(d["intro"], 15, MUTED)], ls=1.22)]); hy += 0.34 * il + 0.2
    hi = d.get("highlight", ""); hib = max(1, math.ceil(len(hi) / 118)) if hi else 0
    hbar_h = (0.32 * hib + 0.3) if hi else 0; hbar_y = 6.86 - hbar_h; cards_top = hy + 0.05
    cards_bot = (hbar_y - 0.22) if hi else 6.62; ch_in = cards_bot - cards_top
    cols = d.get("columns", []); gap = 0.4; total = 11.33; cw = (total - gap) / 2; inner_w = cw - 0.8
    for idx, c in enumerate(cols[:2]):
        cah = c.get("accent", ah); cac = C(cah); cx = x + Inches(idx * (cw + gap))
        rect(s, cx, Inches(cards_top), Inches(cw), Inches(ch_in), fill=WHITE, line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        rect(s, cx, Inches(cards_top), Inches(cw), Inches(0.5), fill=cac, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, cx, Inches(cards_top + 0.25), Inches(cw), Inches(0.25), fill=cac)
        txt(s, cx + Inches(0.32), Inches(cards_top), Inches(cw - 0.6), Inches(0.5), [P([R(c.get("title", ""), 14.5, WHITE, True, font=FONT_SB)], ls=1.0)], anchor=MSO_ANCHOR.MIDDLE)
        by = cards_top + 0.72
        for pt in c.get("points", []):
            nl = max(1, math.ceil(len(pt) / max(1, int((inner_w - 0.3) / ((13 * 0.52) / 72.0)))))
            rect(s, cx + Inches(0.34), Inches(by + 0.07), Inches(0.1), Inches(0.1), fill=cac, shape=MSO_SHAPE.OVAL)
            txt(s, cx + Inches(0.58), Inches(by), Inches(cw - 0.9), Inches(0.24 * nl + 0.1), [P([R(pt, 13, INK_SOFT)], ls=1.16)], anchor=MSO_ANCHOR.TOP); by += 0.225 * nl + 0.16
    if hi:
        rect(s, x, Inches(hbar_y), Inches(11.33), Inches(hbar_h), fill=lighten(ah, 0.86), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, Inches(hbar_y), Inches(0.09), Inches(hbar_h), fill=accent)
        txt(s, x + Inches(0.32), Inches(hbar_y), Inches(10.7), Inches(hbar_h), [P([R(hi, 13.5, darken(ah, 0.12), True)], ls=1.18)], anchor=MSO_ANCHOR.MIDDLE)

def render_processflow(s, d, i):
    bg(s, PAPER); brandbar(s); x = Inches(1.0)
    if d.get("kicker"): kicker(s, x, Inches(1.05), d["kicker"])
    txt(s, x, Inches(1.52), Inches(11.6), Inches(0.8), [P([R(d.get("headline", ""), 34, INK, True, font=FONT_SB)])])
    steps = d.get("steps", []); n = max(1, len(steps)); aw = 0.28; total = 11.33; cw = (total - (n - 1) * aw) / n
    cy = 2.9; chh = 2.95; cxp = 1.0; pal = [C("#0078D4"), C("#8661C5"), C("#D83B01"), C("#107C10"), C("#0E5AA7"), C("#5C2D91")]
    for idx, st in enumerate(steps):
        acc = pal[idx % len(pal)]
        rect(s, Inches(cxp), Inches(cy), Inches(cw), Inches(chh), fill=WHITE, line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        rect(s, Inches(cxp), Inches(cy), Inches(cw), Inches(0.09), fill=acc, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, Inches(cxp + 0.28), Inches(cy + 0.3), Inches(0.5), Inches(0.5), fill=acc, shape=MSO_SHAPE.OVAL)
        txt(s, Inches(cxp + 0.28), Inches(cy + 0.3), Inches(0.5), Inches(0.5), [P([R(str(idx + 1), 19, WHITE, True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(cxp + 0.26), Inches(cy + 0.92), Inches(cw - 0.48), Inches(0.55), [P([R(st.get("title", ""), 15, INK, True, font=FONT_SB)], ls=1.02)])
        txt(s, Inches(cxp + 0.26), Inches(cy + 1.46), Inches(cw - 0.48), Inches(1.35), [P([R(st.get("desc", ""), 10.5, MUTED)], ls=1.18)])
        cxp += cw
        if idx < n - 1:
            txt(s, Inches(cxp), Inches(cy + chh / 2 - 0.28), Inches(aw), Inches(0.55), [P([R("›", 27, C("#BFBFBF"), True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE); cxp += aw
    hs = d.get("highlights", [])
    if hs:
        hy = 6.15; hgap = 0.4; hw = (11.33 - hgap) / 2
        for j, ht in enumerate(hs[:2]):
            hx = 1.0 + j * (hw + hgap)
            rect(s, Inches(hx), Inches(hy), Inches(hw), Inches(0.75), fill=RGBColor(0xF3, 0xF8, 0xFE), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            rect(s, Inches(hx), Inches(hy), Inches(0.08), Inches(0.75), fill=BLUE)
            txt(s, Inches(hx + 0.28), Inches(hy), Inches(hw - 0.52), Inches(0.75), [P([R(ht, 12.5, INK_SOFT, True)], ls=1.15)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, SW - Inches(4.5), SH - Inches(0.62), Inches(4.0), Inches(0.35), [P([R(d.get("footer", "Confidential"), 10.5, MUTED, spacing=0.5)], align=PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)


def _render_dict():
    return {
        "title": render_title, "statement": render_statement, "bullets": render_bullets, "steps": render_steps,
        "columns": render_columns, "quote": render_quote, "cta": render_cta, "feature": render_feature,
        "hero": render_hero, "ecosystem": render_ecosystem, "flow": render_flow,
        "content2col": render_content2col, "processflow": render_processflow,
    }

def _build_deck(spec, out_path):
    global prs, SW, SH, BLANK, WORDMARK, FOOT_LEFT, FOOT_RIGHT, PAGE_LABEL
    deck = spec.get("deck", {}) or {}
    WORDMARK = deck.get("wordmark", "RAPP · Rapid Agent Prototype Platform")
    FOOT_LEFT = deck.get("footLeft", "Microsoft · MCAPS")
    FOOT_RIGHT = deck.get("footRight", "Business Overview · v1")
    PAGE_LABEL = deck.get("pageLabel", "Business Overview")
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]; SW, SH = prs.slide_width, prs.slide_height
    R_ = _render_dict()
    for idx, sd in enumerate(spec.get("slides", []), start=1):
        s = slide(); layout = sd.get("layout", "bullets")
        if layout == "title":
            merged = dict(deck); merged.update(sd); render_title(s, merged, idx)
        else:
            R_.get(layout, render_bullets)(s, sd, idx)
    prs.save(out_path)
    return len(spec.get("slides", []))

def _default_path(spec):
    deck = spec.get("deck", {}) or {}; slides = spec.get("slides", []) or [{}]
    title = deck.get("title") or slides[0].get("headline") or slides[0].get("title") or "RAPP_Deck"
    base = re.sub(r"[^A-Za-z0-9]+", "_", str(title)).strip("_")[:48] or "RAPP_Deck"
    dl = os.path.expanduser("~/Downloads")
    return os.path.join(dl if os.path.isdir(dl) else os.getcwd(), base + ".pptx")


class MicrosoftDeckStudio(BasicAgent):
    def __init__(self):
        self.name = "MicrosoftDeckStudio"
        self.metadata = {
            "name": self.name,
            "description": ("Generate a polished, Microsoft-Fluent-styled PowerPoint (.pptx) from a JSON spec of slides. "
                            "Supported layouts: title, hero, statement, bullets, steps, columns, content2col, feature, "
                            "ecosystem, flow, processflow, quote, cta. Returns the saved .pptx path."),
            "parameters": {
                "type": "object",
                "properties": {
                    "spec": {"type": "string", "description": "Deck spec as a JSON object or JSON string: {deck:{wordmark,footLeft,pageLabel}, slides:[{layout, ...}]}. See the agent docstring for each layout's fields; {h}...{/h} highlights a headline span."},
                    "output_path": {"type": "string", "description": "Absolute .pptx path to write. Optional; defaults to ~/Downloads/<title>.pptx (or CWD)."},
                },
                "required": ["spec"],
            },
        }
        super().__init__(self.name, self.metadata)

    def perform(self, **kwargs):
        spec = kwargs.get("spec")
        if spec is None:
            return "[MicrosoftDeckStudio] 'spec' is required — a JSON object (or string) with a 'slides' array."
        if isinstance(spec, str):
            try:
                spec = json.loads(spec)
            except Exception as e:
                return "[MicrosoftDeckStudio] could not parse 'spec' JSON: %s" % e
        if not isinstance(spec, dict) or not spec.get("slides"):
            return "[MicrosoftDeckStudio] 'spec' must be an object with a non-empty 'slides' array."
        ok, msg = _ensure_engine()
        if not ok:
            return msg
        out = kwargs.get("output_path") or _default_path(spec)
        try:
            out = os.path.abspath(os.path.expanduser(out))
            d = os.path.dirname(out)
            if d and not os.path.isdir(d):
                os.makedirs(d, exist_ok=True)
            n = _build_deck(spec, out)
        except Exception as e:
            return "[MicrosoftDeckStudio] render failed: %s: %s" % (type(e).__name__, e)
        return "[MicrosoftDeckStudio] Generated a %d-slide Microsoft-styled deck -> %s" % (n, out)